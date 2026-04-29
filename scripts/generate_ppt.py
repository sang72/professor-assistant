#!/usr/bin/env python3
"""
Professor Assistant — PPT Generator (Standard 4:3)
Converts a session.md lecture file into a .pptx presentation.

Slide structure:
  Slide 1      : TITLE   — course name, week/topic, professor name
  Slide 2      : TOC     — today's agenda (numbered list)
  Slides 3–N-2 : IMAGE / KEY_POINT / SECTION  (content slides)
  Slide N-1    : QA      — Q&A slide
  Slide N      : END     — identical to slide 1 (title repeated)

Speaker notes = full word-for-word lecture script.

Usage:
  python3 scripts/generate_ppt.py courses/[FOLDER]/lectures/week01/session1.md
"""

import sys, re
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
except ImportError:
    print("Error: python-pptx is not installed.  Run: pip install python-pptx")
    sys.exit(1)

# ── 4:3 canvas ────────────────────────────────────────────────────────────────
W = Inches(10)
H = Inches(7.5)
M = Inches(0.4)   # margin

# ── Colour palette ────────────────────────────────────────────────────────────
DARK_BLUE  = RGBColor(0x1B, 0x35, 0x6E)
ACCENT     = RGBColor(0x2E, 0x6D, 0xB4)
GOLD       = RGBColor(0xC9, 0xA0, 0x2B)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF2, 0xF4, 0xF8)
DARK_GRAY  = RGBColor(0x33, 0x33, 0x33)
MID_GRAY   = RGBColor(0x77, 0x77, 0x77)
IMG_BG     = RGBColor(0xE3, 0xEA, 0xF4)
IMG_BORDER = RGBColor(0x2E, 0x6D, 0xB4)


# ── Primitives ────────────────────────────────────────────────────────────────

def _blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

def _bg(slide, color):
    f = slide.background.fill
    f.solid()
    f.fore_color.rgb = color

def _rect(slide, x, y, w, h, fill, line_color=None, line_pt=0):
    s = slide.shapes.add_shape(1, x, y, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if line_color:
        s.line.color.rgb = line_color
        s.line.width = Pt(line_pt)
    else:
        s.line.fill.background()
    return s

def _txt(slide, x, y, w, h, text, size, bold=False,
         color=DARK_GRAY, align=PP_ALIGN.LEFT, italic=False):
    b = slide.shapes.add_textbox(x, y, w, h)
    tf = b.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    return b

def _notes(slide, text):
    if not text:
        return
    nf = slide.notes_slide.notes_text_frame
    nf.clear()
    for i, line in enumerate(text.split('\n')):
        p = nf.paragraphs[0] if i == 0 else nf.add_paragraph()
        p.text = line


# ── Slide builders ────────────────────────────────────────────────────────────

def build_title(prs, course, week_session, topic, professor, notes=""):
    """Slide 1 and last slide: course + week/topic + professor."""
    slide = _blank(prs)
    _bg(slide, DARK_BLUE)

    # Gold bar top
    _rect(slide, 0, 0, W, Inches(0.12), GOLD)
    # Gold bar bottom
    _rect(slide, 0, H - Inches(0.12), W, Inches(0.12), GOLD)
    # Left accent stripe
    _rect(slide, 0, Inches(0.12), Inches(0.08), H - Inches(0.24), GOLD)

    # Course name
    _txt(slide, Inches(0.35), Inches(1.5), W - Inches(0.55), Inches(1.1),
         course, 32, bold=True, color=WHITE, align=PP_ALIGN.LEFT)

    # Separator line via thin rect
    _rect(slide, Inches(0.35), Inches(2.7), W - Inches(0.75), Inches(0.04), GOLD)

    # Week / Session / Topic
    _txt(slide, Inches(0.35), Inches(2.85), W - Inches(0.55), Inches(0.7),
         week_session, 18, color=GOLD, align=PP_ALIGN.LEFT)
    _txt(slide, Inches(0.35), Inches(3.5), W - Inches(0.55), Inches(1.2),
         topic, 26, bold=True, color=WHITE, align=PP_ALIGN.LEFT)

    # Professor name (bottom-right)
    _txt(slide, Inches(0.35), Inches(6.2), W - Inches(0.75), Inches(0.6),
         professor, 16, color=RGBColor(0xAA, 0xBB, 0xDD), align=PP_ALIGN.RIGHT)

    _notes(slide, notes)
    return slide


def build_toc(prs, items, notes=""):
    """Slide 2: today's agenda as a numbered list."""
    slide = _blank(prs)
    _bg(slide, LIGHT_GRAY)

    # Header bar
    _rect(slide, 0, 0, W, Inches(1.1), DARK_BLUE)
    _rect(slide, 0, Inches(1.1), W, Inches(0.06), GOLD)

    lang_hint = items[0] if items else ""
    header = "오늘의 수업 목차" if _is_korean(lang_hint) else "Today's Agenda"
    _txt(slide, M, Inches(0.15), W - M*2, Inches(0.85),
         header, 26, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # Numbered items
    content_top = Inches(1.4)
    row_h = (H - content_top - Inches(0.3)) / max(len(items), 1)
    for i, item in enumerate(items):
        y = content_top + row_h * i
        # Number circle area
        _txt(slide, Inches(0.5), y, Inches(0.7), row_h,
             str(i + 1), 22, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
        _txt(slide, Inches(1.3), y + Inches(0.05), W - Inches(1.7), row_h,
             item, 20, color=DARK_GRAY)

    _notes(slide, notes)
    return slide


def build_image(prs, title, image_desc, caption="", notes=""):
    """Image-focused content slide — large placeholder, small key phrase."""
    slide = _blank(prs)
    _bg(slide, WHITE)

    # Thin title bar
    _rect(slide, 0, 0, W, Inches(0.9), DARK_BLUE)
    _rect(slide, 0, Inches(0.9), W, Inches(0.05), GOLD)
    _txt(slide, M, Inches(0.1), W - M*2, Inches(0.75),
         title, 22, bold=True, color=WHITE)

    # Large image placeholder (fills ~70% of slide height)
    img_h = Inches(4.8)
    img_y = Inches(1.05)
    _rect(slide, M, img_y, W - M*2, img_h, IMG_BG,
          line_color=IMG_BORDER, line_pt=1.2)

    # Camera icon text + description inside placeholder
    _txt(slide, M + Inches(0.3), img_y + Inches(0.3),
         W - M*2 - Inches(0.6), Inches(0.5),
         "📷", 28, align=PP_ALIGN.CENTER, color=ACCENT)
    _txt(slide, M + Inches(0.3), img_y + Inches(1.0),
         W - M*2 - Inches(0.6), img_h - Inches(1.3),
         image_desc, 15, italic=True, color=ACCENT, align=PP_ALIGN.CENTER)

    # Caption / key phrase at bottom
    if caption:
        _rect(slide, 0, H - Inches(1.15), W, Inches(1.15), DARK_BLUE)
        _txt(slide, M, H - Inches(1.05), W - M*2, Inches(0.9),
             caption, 19, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    _notes(slide, notes)
    return slide


def build_key_point(prs, title, bullets, notes=""):
    """Text-based content slide for definitions, steps, comparisons."""
    slide = _blank(prs)
    _bg(slide, LIGHT_GRAY)

    _rect(slide, 0, 0, W, Inches(1.05), DARK_BLUE)
    _rect(slide, 0, Inches(1.05), W, Inches(0.06), GOLD)
    _txt(slide, M, Inches(0.13), W - M*2, Inches(0.85),
         title, 24, bold=True, color=WHITE)

    tb = slide.shapes.add_textbox(M, Inches(1.25), W - M*2, Inches(5.9))
    tf = tb.text_frame
    tf.word_wrap = True

    for idx, raw in enumerate(bullets):
        stripped = raw.lstrip()
        is_sub   = raw.startswith('    ') or raw.startswith('\t\t')
        ch       = '◦' if is_sub else '●'
        clean    = re.sub(r'^[-*•◦●]\s*', '', stripped)
        text     = f"   {ch}  {clean}" if is_sub else f"{ch}  {clean}"

        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        r = p.add_run()
        r.text = text
        r.font.size  = Pt(16 if is_sub else 20)
        r.font.color.rgb = MID_GRAY if is_sub else DARK_GRAY
        p.space_after = Pt(5 if is_sub else 10)

    _notes(slide, notes)
    return slide


def build_section(prs, title, notes=""):
    """Section-break divider slide."""
    slide = _blank(prs)
    _bg(slide, ACCENT)
    _rect(slide, 0, Inches(3.25), Inches(0.1), Inches(1.0), GOLD)
    _txt(slide, Inches(0.35), Inches(3.0), W - Inches(0.75), Inches(1.5),
         title, 32, bold=True, color=WHITE)
    _notes(slide, notes)
    return slide


def build_qa(prs, notes=""):
    """Q&A slide — dark blue, large Q&A text."""
    slide = _blank(prs)
    _bg(slide, DARK_BLUE)
    _rect(slide, 0, 0, W, Inches(0.12), GOLD)
    _rect(slide, 0, H - Inches(0.12), W, Inches(0.12), GOLD)
    _txt(slide, 0, Inches(2.5), W, Inches(1.5),
         "Q & A", 54, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    _txt(slide, 0, Inches(4.3), W, Inches(0.7),
         "질문 있으신 분?  /  Any questions?",
         18, color=RGBColor(0xAA, 0xBB, 0xDD), align=PP_ALIGN.CENTER)
    _notes(slide, notes)
    return slide


# ── Parser ────────────────────────────────────────────────────────────────────

def _is_korean(text):
    return bool(re.search(r'[가-힣]', text))


def parse_deck(md_text):
    m = re.search(
        r'(?:SECTION\s*3|SLIDE\s*DECK)[^\n]*\n(.*)',
        md_text, re.IGNORECASE | re.DOTALL
    )
    if not m:
        return []

    raw_blocks = re.split(r'\[SLIDE\s*\d+\s*[|｜]\s*', m.group(1))
    slides = []

    for block in raw_blocks:
        block = block.strip()
        if not block:
            continue

        lines      = block.split('\n')
        slide_type = lines[0].strip().rstrip(']').upper()

        s = dict(type=slide_type, title='', subtitle='', professor='',
                 week_session='', topic='', items=[], bullets=[],
                 image_desc='', caption='', notes='')

        in_notes = False
        note_buf = []

        for line in lines[1:]:
            if in_notes:
                note_buf.append(line)
                continue
            if   line.startswith('Title:'):       s['title']        = line[6:].strip()
            elif line.startswith('Subtitle:'):    s['subtitle']     = line[9:].strip()
            elif line.startswith('Professor:'):   s['professor']    = line[10:].strip()
            elif line.startswith('WeekSession:'): s['week_session'] = line[12:].strip()
            elif line.startswith('Topic:'):       s['topic']        = line[6:].strip()
            elif line.startswith('Image:'):       s['image_desc']   = line[6:].strip()
            elif line.startswith('Caption:'):     s['caption']      = line[8:].strip()
            elif line.startswith('- ') or re.match(r'\s{2,}-', line):
                s['bullets'].append(line)
                if slide_type in ('TOC', 'AGENDA'):
                    clean = re.sub(r'^\s*-\s*', '', line).strip()
                    s['items'].append(clean)
            elif line.startswith('Notes:'):
                in_notes = True
                rest = line[6:].strip()
                if rest:
                    note_buf.append(rest)

        s['notes'] = '\n'.join(note_buf).strip()
        slides.append(s)

    return slides


def fallback(md_text):
    """Build basic slide list from lesson-plan header when no SLIDE DECK section exists."""
    def _get(pattern):
        m = re.search(pattern, md_text)
        return m.group(1).strip() if m else ''

    course    = _get(r'\|\s*Course\s*\|\s*(.+?)\s*\|')
    week      = _get(r'\|\s*Week\s*\|\s*(.+?)\s*\|')
    session   = _get(r'\|\s*Session\s*\|\s*(.+?)\s*\|')
    topic     = _get(r'\|\s*Topic\s*\|\s*(.+?)\s*\|')
    professor = _get(r'Professor[:\s]+([^\n|]+)')

    slides = [
        dict(type='TITLE', title=course or 'Lecture',
             week_session=f"Week {week} · Session {session}",
             topic=topic, professor=professor,
             subtitle='', items=[], bullets=[], image_desc='', caption='', notes=''),
        dict(type='TOC', items=['Topic 1', 'Topic 2', 'Topic 3'],
             title='', subtitle='', professor='', week_session='', topic='',
             bullets=[], image_desc='', caption='', notes=''),
    ]

    rows = re.findall(r'\|\s*[\d:]+\s*\|\s*[\d ]+\s*min\s*\|\s*([^|]{3,}?)\s*\|', md_text)
    for act in rows:
        a = act.strip()
        if a and a.lower() not in ('activity', 'duration'):
            slides.append(dict(type='KEY_POINT', title=a, bullets=[],
                               subtitle='', professor='', week_session='', topic='',
                               items=[], image_desc='', caption='', notes=''))

    slides.append(dict(type='QA', title='', subtitle='', professor='',
                       week_session='', topic='', items=[], bullets=[],
                       image_desc='', caption='', notes=''))
    slides.append(dict(type='END', title=course or 'Lecture',
                       week_session=f"Week {week} · Session {session}",
                       topic=topic, professor=professor,
                       subtitle='', items=[], bullets=[], image_desc='', caption='', notes=''))
    return slides


# ── Main ──────────────────────────────────────────────────────────────────────

def generate(md_path_str):
    path = Path(md_path_str)
    if not path.exists():
        print(f"Error: file not found — {path}")
        sys.exit(1)

    text   = path.read_text(encoding='utf-8')
    slides = parse_deck(text)

    if not slides:
        print("No SLIDE DECK section found — building fallback from lesson plan...")
        slides = fallback(text)

    # Pull title-slide data for END slide reuse
    title_s = next((s for s in slides if s['type'] in ('TITLE', 'END')), None)

    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H

    for s in slides:
        t = s['type']

        if t in ('TITLE',):
            build_title(prs,
                        s['title'],
                        s.get('week_session', s.get('subtitle', '')),
                        s.get('topic', ''),
                        s.get('professor', ''),
                        s['notes'])

        elif t in ('TOC', 'AGENDA', 'CONTENTS'):
            items = s['items'] if s['items'] else [b.lstrip('- ').strip() for b in s['bullets']]
            build_toc(prs, items, s['notes'])

        elif t in ('IMAGE', 'PHOTO', 'VISUAL'):
            build_image(prs, s['title'], s['image_desc'], s['caption'], s['notes'])

        elif t in ('SECTION', 'DIVIDER'):
            build_section(prs, s['title'], s['notes'])

        elif t in ('QA', 'Q&A', 'QUESTIONS'):
            build_qa(prs, s['notes'])

        elif t in ('END', 'CLOSING', 'FINAL'):
            # Repeat the title slide
            ref = title_s or s
            build_title(prs,
                        ref['title'],
                        ref.get('week_session', ref.get('subtitle', '')),
                        ref.get('topic', ''),
                        ref.get('professor', ''),
                        s['notes'])

        else:  # KEY_POINT, CONTENT, default
            build_key_point(prs, s['title'], s['bullets'], s['notes'])

    out = path.parent / (path.stem + '_slides.pptx')
    prs.save(out)
    print(f"✓ PPT 저장 완료 → {out}")
    return out


if __name__ == '__main__':
    if len(sys.argv) < 2:
        print("사용법: python3 scripts/generate_ppt.py <경로/session.md>")
        print("예시:  python3 scripts/generate_ppt.py courses/ENG101/lectures/week01/session1.md")
        sys.exit(1)
    generate(sys.argv[1])
